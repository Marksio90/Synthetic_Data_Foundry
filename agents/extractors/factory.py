"""
agents/extractors/factory.py — Multi-format content extractor for Gap Scout.

Supported formats (all optional deps degrade gracefully to PlainTextExtractor):
  pdf   — pdfminer.six  + pymupdf fallback
  html  — trafilatura   + bs4/lxml fallback + regex last-resort
  docx  — python-docx   (text + tables)
  xlsx  — openpyxl      → JSON schema + 10-row sample
  pptx  — python-pptx   → slide-by-slide text + speaker notes
  srt   — built-in regex (no extra dep)
  vtt   — built-in regex (no extra dep)
  xml   — lxml          + regex fallback
  csv   — stdlib csv    (no extra dep)
  txt / md / json  — PlainTextExtractor (UTF-8 decode)

Audio/video extractors (mp3, mp4, wav) are wired in Step 8 (Whisper).

Usage:
    from agents.extractors.factory import ExtractorFactory

    text = ExtractorFactory.create("pdf").extract(raw_bytes)
    text = ExtractorFactory.from_path("report.xlsx").extract(raw_bytes)
    snippet = text[:500]  # store only 500-char snippets per copyright policy
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Union

logger = logging.getLogger(__name__)

ContentInput = Union[bytes, str]


# ===========================================================================
# Base
# ===========================================================================


class BaseExtractor(ABC):
    """
    Abstract extractor interface.
    extract() always returns a str — empty string on total failure.
    """

    @abstractmethod
    def extract(self, content: ContentInput) -> str:
        ...

    @staticmethod
    def _to_bytes(content: ContentInput) -> bytes:
        return content if isinstance(content, bytes) else content.encode("utf-8", errors="replace")

    @staticmethod
    def _to_str(content: ContentInput, encoding: str = "utf-8") -> str:
        return content if isinstance(content, str) else content.decode(encoding, errors="replace")


# ===========================================================================
# Plain text (fallback for everything)
# ===========================================================================


class PlainTextExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        return self._to_str(content)


# ===========================================================================
# PDF — pdfminer.six primary, pymupdf fallback
# ===========================================================================


class PDFExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_bytes(content)

        # Primary: pdfminer.six
        try:
            from pdfminer.high_level import extract_text_to_fp  # type: ignore
            from pdfminer.layout import LAParams  # type: ignore

            output = io.StringIO()
            extract_text_to_fp(io.BytesIO(raw), output, laparams=LAParams())
            text = output.getvalue().strip()
            if text:
                return text
        except ImportError:
            pass
        except Exception as exc:
            logger.debug("pdfminer failed for PDF: %s — trying pymupdf", exc)

        # Fallback: pymupdf (fitz)
        try:
            import fitz  # type: ignore  # pymupdf

            doc = fitz.open(stream=raw, filetype="pdf")
            pages = [page.get_text() for page in doc]
            doc.close()
            text = "\n".join(pages).strip()
            if text:
                return text
        except ImportError:
            pass
        except Exception as exc:
            logger.debug("pymupdf failed for PDF: %s", exc)

        logger.warning(
            "No PDF extractor available — install pdfminer.six or pymupdf"
        )
        return PlainTextExtractor().extract(raw)


# ===========================================================================
# HTML — trafilatura primary, bs4/lxml fallback, regex last resort
# ===========================================================================


class HTMLExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw_str = self._to_str(content)

        # Primary: trafilatura (article-optimised, respects boilerplate removal)
        try:
            import trafilatura  # type: ignore

            text = trafilatura.extract(
                raw_str,
                include_comments=False,
                include_tables=True,
                no_fallback=False,
            )
            if text:
                return text.strip()
        except ImportError:
            pass
        except Exception as exc:
            logger.debug("trafilatura failed: %s — trying bs4", exc)

        # First fallback: bs4 + lxml
        try:
            from bs4 import BeautifulSoup  # type: ignore

            parser = "lxml"
            try:
                import lxml  # type: ignore  # noqa: F401
            except ImportError:
                parser = "html.parser"

            soup = BeautifulSoup(raw_str, parser)
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            pass
        except Exception as exc:
            logger.debug("bs4 failed: %s — using regex fallback", exc)

        # Last resort: strip tags with regex
        return re.sub(r"<[^>]+>", " ", raw_str).strip()


# ===========================================================================
# DOCX — python-docx (text + tables)
# ===========================================================================


class DocxExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_bytes(content)
        try:
            import docx  # type: ignore

            doc = docx.Document(io.BytesIO(raw))
            parts: list[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join(c.text.strip() for c in row.cells)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
        except ImportError:
            logger.warning("python-docx not installed — install python-docx")
        except Exception as exc:
            logger.debug("DOCX extraction failed: %s", exc)
        return PlainTextExtractor().extract(raw)


# ===========================================================================
# XLSX — openpyxl → JSON schema + 10-row sample per sheet
# ===========================================================================


class XlsxExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_bytes(content)
        try:
            import openpyxl  # type: ignore

            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            sheets: dict[str, dict] = {}
            for name in wb.sheetnames:
                ws = wb[name]
                all_rows = list(ws.iter_rows(values_only=True))
                if not all_rows:
                    continue
                headers = [str(c) if c is not None else f"col_{i}" for i, c in enumerate(all_rows[0])]
                sample = [
                    {h: (str(v) if v is not None else "") for h, v in zip(headers, row)}
                    for row in all_rows[1:11]
                ]
                sheets[name] = {
                    "columns": headers,
                    "row_count": ws.max_row,
                    "sample_rows": sample,
                }
            wb.close()
            return json.dumps(sheets, ensure_ascii=False, indent=2)
        except ImportError:
            logger.warning("openpyxl not installed — install openpyxl")
        except Exception as exc:
            logger.debug("XLSX extraction failed: %s", exc)
        return PlainTextExtractor().extract(raw)


# ===========================================================================
# PPTX — python-pptx, slide-by-slide text + speaker notes
# ===========================================================================


class PptxExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_bytes(content)
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(io.BytesIO(raw))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text_frame.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text_frame") and shape.text_frame.text.strip()
                ]
                parts.append(f"[Slide {i}] " + " | ".join(texts))
                if slide.has_notes_slide:
                    ntf = slide.notes_slide.notes_text_frame
                    if ntf and ntf.text.strip():
                        parts.append(f"  Notes: {ntf.text.strip()}")
            return "\n".join(parts)
        except ImportError:
            logger.warning("python-pptx not installed — install python-pptx")
        except Exception as exc:
            logger.debug("PPTX extraction failed: %s", exc)
        return PlainTextExtractor().extract(raw)


# ===========================================================================
# SRT / VTT — built-in regex, no extra dep
# ===========================================================================


class SubtitleExtractor(BaseExtractor):
    _TIMESTAMP_RE = re.compile(
        r"\d{2}:\d{2}:\d{2}[,.]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[,.]\d{3}"
    )
    _SEQUENCE_RE = re.compile(r"^\d+\s*$", re.MULTILINE)
    _WEBVTT_HEADER = re.compile(r"^WEBVTT.*$", re.MULTILINE)
    _INLINE_TAG = re.compile(r"<[^>]+>")

    def extract(self, content: ContentInput) -> str:
        raw = self._to_str(content)
        raw = self._WEBVTT_HEADER.sub("", raw)
        raw = self._TIMESTAMP_RE.sub("", raw)
        raw = self._SEQUENCE_RE.sub("", raw)
        raw = self._INLINE_TAG.sub("", raw)
        lines = [line.strip() for line in raw.splitlines() if line.strip()]
        return " ".join(lines)


# ===========================================================================
# XML — lxml itertext with regex fallback
# ===========================================================================


class XMLExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_bytes(content)
        try:
            from lxml import etree  # type: ignore

            root = etree.fromstring(raw)
            texts = [t.strip() for t in root.itertext() if t.strip()]
            return "\n".join(texts)
        except ImportError:
            pass
        except Exception as exc:
            logger.debug("lxml XML extraction failed: %s", exc)
        decoded = raw.decode("utf-8", errors="replace")
        return re.sub(r"<[^>]+>", " ", decoded).strip()


# ===========================================================================
# CSV — stdlib only
# ===========================================================================


class CSVExtractor(BaseExtractor):
    def extract(self, content: ContentInput) -> str:
        raw = self._to_str(content)
        try:
            reader = csv.reader(io.StringIO(raw))
            rows = list(reader)
            if not rows:
                return ""
            header = ", ".join(rows[0])
            body = "\n".join(", ".join(r) for r in rows[1:21])
            return f"Columns: {header}\n{body}"
        except Exception as exc:
            logger.debug("CSV extraction failed: %s", exc)
        return raw


# ===========================================================================
# Audio / Video — faster-whisper primary, Replicate API fallback (Step 8)
# ===========================================================================


class AudioExtractor(BaseExtractor):
    """
    Transcribes audio (mp3, wav, m4a, ogg, flac) and video (mp4, mkv, webm…)
    using faster-whisper locally or Replicate API as fallback.
    See agents/extractors/audio.py for full implementation details.
    """

    def extract(self, content: ContentInput) -> str:
        from agents.extractors.audio import WhisperAudioExtractor
        return WhisperAudioExtractor().extract(
            content if isinstance(content, bytes) else content.encode("utf-8", errors="replace")
        )


# ===========================================================================
# Factory
# ===========================================================================

_FORMAT_MAP: dict[str, type[BaseExtractor]] = {
    "pdf":  PDFExtractor,
    "html": HTMLExtractor,
    "htm":  HTMLExtractor,
    "docx": DocxExtractor,
    "xlsx": XlsxExtractor,
    "xls":  XlsxExtractor,
    "pptx": PptxExtractor,
    "ppt":  PptxExtractor,
    "srt":  SubtitleExtractor,
    "vtt":  SubtitleExtractor,
    "xml":  XMLExtractor,
    "csv":  CSVExtractor,
    "txt":  PlainTextExtractor,
    "md":   PlainTextExtractor,
    "json": PlainTextExtractor,
    "jsonld": PlainTextExtractor,
    # Audio formats
    "mp3":  AudioExtractor,
    "wav":  AudioExtractor,
    "m4a":  AudioExtractor,
    "ogg":  AudioExtractor,
    "flac": AudioExtractor,
    "aac":  AudioExtractor,
    "opus": AudioExtractor,
    "wma":  AudioExtractor,
    # Video formats (audio track extracted via ffmpeg before transcription)
    "mp4":  AudioExtractor,
    "mkv":  AudioExtractor,
    "webm": AudioExtractor,
    "mov":  AudioExtractor,
    "avi":  AudioExtractor,
    "m4v":  AudioExtractor,
    "wmv":  AudioExtractor,
}


class ExtractorFactory:
    """
    Factory that selects the appropriate BaseExtractor for a given format.

    Usage:
        text = ExtractorFactory.create("pdf").extract(raw_bytes)
        text = ExtractorFactory.from_path("/tmp/report.xlsx").extract(raw_bytes)
        snippet = text[:500]   # copyright policy: store ≤ 500-char snippets only
    """

    @staticmethod
    def create(fmt: str) -> BaseExtractor:
        """Return extractor for an explicit format string (e.g. 'pdf', 'html')."""
        key = fmt.lower().lstrip(".")
        cls = _FORMAT_MAP.get(key, PlainTextExtractor)
        return cls()

    @staticmethod
    def from_path(path: Union[str, Path]) -> BaseExtractor:
        """Infer extractor from file extension."""
        suffix = Path(path).suffix.lower().lstrip(".")
        return ExtractorFactory.create(suffix)

    @staticmethod
    def supported_formats() -> list[str]:
        """Return sorted list of supported format strings."""
        return sorted(_FORMAT_MAP.keys())
